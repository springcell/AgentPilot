"""
file_ops.py — File operations module
AI invokes via JSON: list, find, read, write, write_chunk, append, patch, insert,
delete_lines, mkdir, copy, move, delete, exists, stat, tree, find_program, launch,
backup, restore, history.

Mutation actions (write / write_chunk / patch / insert / delete_lines / append)
auto-backup before modifying, unless backup=false is passed.
"""

import os
import re
import json
import base64
import shutil
import fnmatch
import urllib.request
import urllib.error
import tempfile
from datetime import datetime
from pathlib import Path


# ── 内部工具 ──────────────────────────────────────────────

def _expand(path: str) -> str:
    """展开环境变量和 ~ 并返回绝对路径。"""
    return str(Path(os.path.expandvars(os.path.expanduser(path))).resolve())


def _read_text(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "gbk", "gb18030", "latin-1"):
        try:
            return Path(path).read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return Path(path).read_bytes().decode("utf-8", errors="replace")


_BINARY_EXTENSIONS = {
    ".pptx", ".ppt", ".xlsx", ".xls", ".docx", ".doc",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".ico", ".tiff",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv",
    ".zip", ".rar", ".7z", ".gz", ".exe", ".dll", ".so", ".pdf",
}


def _is_binary_file(path: str) -> bool:
    ext = os.path.splitext(path)[1].lower()
    if ext in _BINARY_EXTENSIONS:
        return True
    try:
        with open(path, "rb") as f:
            return b"\x00" in f.read(512)
    except Exception:
        return False


def _write_text(path: str, content: str, encoding: str = "utf-8") -> None:
    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding=encoding)


# ── 历史备份 ───────────────────────────────────────────────

_HISTORY_DIR_NAME = ".file_ops_history"
_MAX_BACKUPS_PER_FILE = 20          # 每个文件最多保留的备份数量


def _history_dir(file_path: str) -> Path:
    """返回该文件对应的备份目录（与文件同级的 .file_ops_history/）。"""
    return Path(file_path).parent / _HISTORY_DIR_NAME


def _backup_name(file_path: str) -> str:
    """生成备份文件名：<原文件名>.<时间戳>.bak"""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    return f"{Path(file_path).name}.{ts}.bak"


def _make_backup(file_path: str) -> str | None:
    """
    把 file_path 当前内容备份到 .file_ops_history/ 目录。
    返回备份文件的绝对路径；文件不存在时返回 None（首次写入无需备份）。
    超出 _MAX_BACKUPS_PER_FILE 时自动删除最旧的备份。
    """
    p = Path(file_path)
    if not p.is_file():
        return None

    hdir = _history_dir(file_path)
    hdir.mkdir(parents=True, exist_ok=True)

    bak_path = hdir / _backup_name(file_path)
    shutil.copy2(file_path, bak_path)

    # 清理旧备份（按名称排序，最旧在前）
    prefix = p.name + "."
    all_baks = sorted(hdir.glob(f"{prefix}*.bak"))
    for old in all_baks[:-_MAX_BACKUPS_PER_FILE]:
        try:
            old.unlink()
        except OSError:
            pass

    return str(bak_path)


def _list_backups(file_path: str) -> list[dict]:
    """返回该文件的所有备份，按时间降序（最新在前）。"""
    p = Path(file_path)
    hdir = _history_dir(file_path)
    prefix = p.name + "."
    baks = sorted(hdir.glob(f"{prefix}*.bak"), reverse=True) if hdir.is_dir() else []
    result = []
    for b in baks:
        try:
            stat = b.stat()
            result.append({
                "backup": str(b),
                "name": b.name,
                "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
                "size_bytes": stat.st_size,
            })
        except OSError:
            pass
    return result


# ── action 处理函数 ────────────────────────────────────────

def _op_list(params: dict) -> dict:
    """列出目录内容（非递归）。"""
    path = _expand(params.get("path", "."))
    if not os.path.isdir(path):
        return {"ok": False, "error": f"Not a directory: {path}"}
    entries = []
    for name in sorted(os.listdir(path)):
        full = os.path.join(path, name)
        is_dir = os.path.isdir(full)
        entries.append({
            "name": name,
            "type": "dir" if is_dir else "file",
            "size": 0 if is_dir else os.path.getsize(full),
        })
    return {"ok": True, "path": path, "count": len(entries), "entries": entries}


def _op_find(params: dict) -> dict:
    """递归搜索文件，支持通配符和正则。"""
    root     = _expand(params.get("path", "."))
    pattern  = params.get("pattern", "*")          # 通配符，如 *.txt
    regex    = params.get("regex", None)            # 正则，如 .*新闻.*\.txt
    max_depth = params.get("max_depth", 10)
    max_results = params.get("max_results", 200)
    match_type = params.get("type", "any")          # file / dir / any

    results = []
    re_obj = re.compile(regex) if regex else None

    for dirpath, dirnames, filenames in os.walk(root):
        depth = dirpath.replace(root, "").count(os.sep)
        if depth >= max_depth:
            dirnames.clear()
            continue

        candidates = []
        if match_type in ("file", "any"):
            candidates += [(f, os.path.join(dirpath, f)) for f in filenames]
        if match_type in ("dir", "any"):
            candidates += [(d, os.path.join(dirpath, d)) for d in dirnames]

        for name, full in candidates:
            hit = fnmatch.fnmatch(name, pattern) if not re_obj else bool(re_obj.search(name))
            if hit:
                results.append({"name": name, "path": full,
                                 "type": "dir" if os.path.isdir(full) else "file"})
            if len(results) >= max_results:
                return {"ok": True, "count": len(results),
                        "truncated": True, "results": results}

    return {"ok": True, "count": len(results), "truncated": False, "results": results}


def _op_read(params: dict) -> dict:
    """读取文件内容，支持行范围截取。二进制文件返回元数据摘要。"""
    path = _expand(params.get("path", ""))
    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}

    if _is_binary_file(path):
        ext = os.path.splitext(path)[1].lower()
        size = os.path.getsize(path)
        summary = f"[Binary file: {ext}, size={size} bytes, path={path}]"

        if ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(path)
                texts = []
                for i, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(f"Slide {i}: {shape.text.strip()[:200]}")
                summary = (
                    f"[PPTX: {len(prs.slides)} slides, path={path}]\n"
                    + "\n".join(texts[:30])
                )
            except ImportError:
                summary += "\n(pip install python-pptx to read content)"
            except Exception as e:
                summary += f"\n(Could not parse: {e})"

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                lines = []
                for sheet in wb.sheetnames[:3]:
                    ws = wb[sheet]
                    lines.append(f"Sheet: {sheet}")
                    for row in list(ws.iter_rows(values_only=True))[:10]:
                        lines.append("  " + str(row))
                summary = f"[XLSX: {len(wb.sheetnames)} sheet(s), path={path}]\n" + "\n".join(lines)
            except ImportError:
                summary += "\n(pip install openpyxl to read content)"
            except Exception as e:
                summary += f"\n(Could not parse: {e})"

        elif ext in (".docx", ".doc"):
            try:
                import docx
                doc = docx.Document(path)
                paras = [p.text for p in doc.paragraphs if p.text.strip()][:20]
                summary = f"[DOCX: {len(doc.paragraphs)} paragraphs, path={path}]\n" + "\n".join(paras)
            except ImportError:
                summary += "\n(pip install python-docx to read content)"
            except Exception as e:
                summary += f"\n(Could not parse: {e})"

        return {
            "ok": True,
            "path": path,
            "binary": True,
            "total_lines": 0,
            "shown_lines": "0-0",
            "content": summary,
        }

    content = _read_text(path)
    lines   = content.splitlines()
    start = params.get("line_start", 1) - 1
    end   = params.get("line_end", len(lines))
    snippet = "\n".join(lines[start:end])
    return {
        "ok": True,
        "path": path,
        "binary": False,
        "total_lines": len(lines),
        "shown_lines": f"{start+1}-{min(end, len(lines))}",
        "content": snippet,
    }



def _op_write(params: dict) -> dict:
    """写入/覆盖文件，自动创建父目录。写入前自动备份原文件。"""
    path     = _expand(params.get("path", ""))
    content  = params.get("content", "")
    encoding = params.get("encoding", "utf-8")
    do_backup = params.get("backup", True)

    if not path:
        return {"ok": False, "error": "path cannot be empty"}

    bak = _make_backup(path) if do_backup else None
    _write_text(path, content, encoding)
    result = {"ok": True, "path": path,
              "bytes": len(content.encode(encoding, errors="replace"))}
    if bak:
        result["backup"] = bak
    return result


def _op_write_chunk(params: dict) -> dict:
    """
    分段写入：用 content 替换文件中 line_start..line_end 的内容，其余行保留不变。
    写入前自动备份。

    params:
      path        — 目标文件（不存在则创建，视为空文件）
      content     — 要写入的新内容（替换指定行范围）
      line_start  — 起始行号（1-based，默认 1）
      line_end    — 结束行号（1-based，含；默认替换到文件末尾）
                    传 0 或省略 = 追加到文件末尾（不替换任何行）
      encoding    — 文件编码（默认 utf-8）
      backup      — 是否备份（默认 True）

    典型用法：
      • 只改第 10~20 行：line_start=10, line_end=20, content="新内容"
      • 在文件末尾追加：line_start=0（或不传），line_end=0
      • 整体替换前 N 行：line_start=1, line_end=N
    """
    path      = _expand(params.get("path", ""))
    new_chunk = params.get("content", "")
    encoding  = params.get("encoding", "utf-8")
    do_backup = params.get("backup", True)

    if not path:
        return {"ok": False, "error": "path cannot be empty"}

    # 读取现有内容（文件不存在视为空）
    p = Path(path)
    if p.exists() and not p.is_file():
        return {"ok": False, "error": f"Not a file: {path}"}

    existing = _read_text(path) if p.is_file() else ""
    lines = existing.splitlines(keepends=True)
    total = len(lines)

    line_start = params.get("line_start", 0)
    line_end   = params.get("line_end",   0)

    # line_start=0 或均为 0 → 纯追加模式
    if line_start == 0 and line_end == 0:
        bak = _make_backup(path) if do_backup and p.is_file() else None
        chunk_with_nl = new_chunk if new_chunk.endswith("\n") else new_chunk + "\n"
        final = existing + chunk_with_nl
        _write_text(path, final, encoding)
        result = {"ok": True, "path": path, "mode": "append",
                  "total_lines": len(final.splitlines()),
                  "bytes": len(final.encode(encoding, errors="replace"))}
        if bak:
            result["backup"] = bak
        return result

    # 正常替换模式
    s = max(0, line_start - 1)           # 0-based inclusive start
    e = line_end if line_end else total   # 0-based exclusive end

    # 确保 chunk 以换行结尾（除非是文件最后一块且原文件没有尾换行）
    chunk_lines = new_chunk.splitlines(keepends=True)
    if chunk_lines and not chunk_lines[-1].endswith("\n"):
        chunk_lines[-1] += "\n"

    new_lines = lines[:s] + chunk_lines + lines[e:]
    final = "".join(new_lines)

    bak = _make_backup(path) if do_backup and p.is_file() else None
    _write_text(path, final, encoding)

    result = {
        "ok": True, "path": path, "mode": "chunk",
        "replaced_lines": f"{line_start}-{min(e, total)}",
        "new_chunk_lines": len(chunk_lines),
        "total_lines": len(new_lines),
        "bytes": len(final.encode(encoding, errors="replace")),
    }
    if bak:
        result["backup"] = bak
    return result


def _op_append(params: dict) -> dict:
    """追加内容到文件末尾（不存在则创建）。"""
    path    = _expand(params.get("path", ""))
    content = params.get("content", "")
    newline = params.get("newline", True)
    do_backup = params.get("backup", True)

    if not path:
        return {"ok": False, "error": "path cannot be empty"}

    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)

    existing = p.read_text(encoding="utf-8", errors="replace") if p.exists() else ""
    bak = _make_backup(path) if do_backup and p.is_file() else None

    if newline and existing and not existing.endswith("\n"):
        content = "\n" + content

    with open(path, "a", encoding="utf-8") as f:
        f.write(content)

    result = {"ok": True, "path": path, "appended_bytes": len(content.encode())}
    if bak:
        result["backup"] = bak
    return result


def _op_patch(params: dict) -> dict:
    """
    在文件中替换字符串，支持多处替换。
    replacements: [ {"old": "...", "new": "..."}, ... ]
    use_regex: true 时 old 视为正则表达式
    """
    path  = _expand(params.get("path", ""))
    reps  = params.get("replacements", [])
    use_re = params.get("use_regex", False)
    do_backup = params.get("backup", True)

    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}
    if not reps:
        return {"ok": False, "error": "replacements cannot be empty"}

    content = _read_text(path)
    original = content
    total_count = 0

    for item in reps:
        old = item.get("old", "")
        new = item.get("new", "")
        count_param = item.get("count", 0)

        if use_re:
            flags = re.MULTILINE | (re.IGNORECASE if item.get("ignore_case") else 0)
            new_content, n = re.subn(old, new, content, count=count_param, flags=flags)
        else:
            n = content.count(old) if not count_param else min(content.count(old), count_param)
            new_content = content.replace(old, new, count_param or -1)

        content = new_content
        total_count += n

    if content == original:
        return {"ok": False, "path": path, "replaced": 0,
                "error": "No match; file unchanged. Try file_op read then file_op write_chunk to overwrite."}

    bak = _make_backup(path) if do_backup else None
    _write_text(path, content)
    result = {"ok": True, "path": path, "replaced": total_count}
    if bak:
        result["backup"] = bak
    return result


def _op_insert(params: dict) -> dict:
    """在指定行号前或后插入内容。"""
    path     = _expand(params.get("path", ""))
    line_num = params.get("line", 1)
    content  = params.get("content", "")
    after    = params.get("after", False)
    do_backup = params.get("backup", True)

    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}

    lines = _read_text(path).splitlines(keepends=True)
    idx   = max(0, min(line_num - 1, len(lines)))
    if after:
        idx += 1

    insert_lines = [l + "\n" for l in content.splitlines()]
    lines[idx:idx] = insert_lines

    bak = _make_backup(path) if do_backup else None
    _write_text(path, "".join(lines))
    result = {"ok": True, "path": path, "inserted_at_line": idx + 1,
              "inserted_lines": len(insert_lines)}
    if bak:
        result["backup"] = bak
    return result


def _op_delete_lines(params: dict) -> dict:
    """删除指定行范围（包含两端）。"""
    path  = _expand(params.get("path", ""))
    start = params.get("line_start", 1) - 1
    end   = params.get("line_end", start + 1)
    do_backup = params.get("backup", True)

    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}

    lines = _read_text(path).splitlines(keepends=True)
    deleted = lines[start:end]
    lines[start:end] = []

    bak = _make_backup(path) if do_backup else None
    _write_text(path, "".join(lines))
    result = {"ok": True, "path": path,
              "deleted_lines": len(deleted), "range": f"{start+1}-{end}"}
    if bak:
        result["backup"] = bak
    return result


def _op_mkdir(params: dict) -> dict:
    path = _expand(params.get("path", ""))
    Path(path).mkdir(parents=True, exist_ok=True)
    return {"ok": True, "path": path}


def _op_copy(params: dict) -> dict:
    src  = _expand(params.get("src", ""))
    dst  = _expand(params.get("dst", ""))
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {src}"}
    Path(dst).parent.mkdir(parents=True, exist_ok=True)
    if os.path.isdir(src):
        shutil.copytree(src, dst, dirs_exist_ok=True)
    else:
        shutil.copy2(src, dst)
    return {"ok": True, "src": src, "dst": dst}


def _op_move(params: dict) -> dict:
    src = _expand(params.get("src", ""))
    dst = _expand(params.get("dst", ""))
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {src}"}
    Path(dst).parent.mkdir(parents=True, exist_ok=True)
    shutil.move(src, dst)
    return {"ok": True, "src": src, "dst": dst}


def _op_delete(params: dict) -> dict:
    path  = _expand(params.get("path", ""))
    force = params.get("force", False)
    if not os.path.exists(path):
        return {"ok": False, "error": f"Path not found: {path}"}
    if os.path.isdir(path):
        if force:
            shutil.rmtree(path)
        else:
            return {"ok": False,
                    "error": "Delete directory requires force: true (irreversible)"}
    else:
        os.remove(path)
    return {"ok": True, "deleted": path}


def _op_exists(params: dict) -> dict:
    path = _expand(params.get("path", ""))
    exists = os.path.exists(path)
    return {"ok": True, "path": path, "exists": exists,
            "type": ("dir" if os.path.isdir(path) else "file") if exists else None}


def _op_stat(params: dict) -> dict:
    path = _expand(params.get("path", ""))
    if not os.path.exists(path):
        return {"ok": False, "error": f"Not found: {path}"}
    s = os.stat(path)
    return {
        "ok": True, "path": path,
        "size_bytes": s.st_size,
        "modified": datetime.fromtimestamp(s.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
        "created":  datetime.fromtimestamp(s.st_ctime).strftime("%Y-%m-%d %H:%M:%S"),
        "is_dir":   os.path.isdir(path),
    }


def _op_tree(params: dict) -> dict:
    """输出目录树，返回文本形式。"""
    root      = _expand(params.get("path", "."))
    max_depth = params.get("max_depth", 3)
    lines     = []

    def _walk(path, prefix, depth):
        if depth > max_depth:
            return
        try:
            items = sorted(os.listdir(path))
        except PermissionError:
            return
        for i, name in enumerate(items):
            full  = os.path.join(path, name)
            is_last = (i == len(items) - 1)
            connector = "└── " if is_last else "├── "
            lines.append(prefix + connector + name)
            if os.path.isdir(full):
                extension = "    " if is_last else "│   "
                _walk(full, prefix + extension, depth + 1)

    lines.append(root)
    _walk(root, "", 1)
    return {"ok": True, "tree": "\n".join(lines)}




def _op_find_program(params: dict) -> dict:
    """
    在 Windows 常见位置搜索已安装程序的可执行文件路径。
    params: name (程序名，如 "Photoshop" / "Unity" / "notepad")
    返回所有匹配到的路径列表。
    """
    import glob
    name = params.get("name", "").strip()
    if not name:
        return {"ok": False, "error": "name cannot be empty"}

    # 搜索范围：常见安装目录 + PATH
    search_roots = [
        r"C:\Program Files",
        r"C:\Program Files (x86)",
        os.path.expandvars(r"%LOCALAPPDATA%\Programs"),
        os.path.expandvars(r"%APPDATA%\Microsoft\Windows\Start Menu\Programs"),
        r"D:\Program Files",
        r"D:\Program Files (x86)",
    ]
    # 常见程序名 → 可执行文件名映射
    name_map = {
        "photoshop":  ["Photoshop.exe"],
        "unity":      ["Unity.exe", "UnityHub.exe"],
        "unityhub":   ["UnityHub.exe"],
        "vscode":     ["Code.exe"],
        "vs code":    ["Code.exe"],
        "chrome":     ["chrome.exe"],
        "firefox":    ["firefox.exe"],
        "notepad":    ["notepad.exe"],
        "notepad++":  ["notepad++.exe"],
        "blender":    ["blender.exe"],
        "steam":      ["steam.exe"],
        "obs":        ["obs64.exe", "obs32.exe", "obs.exe"],
    }

    key = name.lower()
    exe_names = name_map.get(key, [f"{name}.exe", f"{name}64.exe"])

    results = []
    seen = set()

    # 1. 在搜索根目录递归查找
    for root in search_roots:
        if not os.path.isdir(root):
            continue
        for exe in exe_names:
            pattern = os.path.join(root, "**", exe)
            for match in glob.glob(pattern, recursive=True):
                norm = os.path.normpath(match)
                if norm not in seen:
                    seen.add(norm)
                    results.append({"path": norm, "source": "search"})

    # 2. 从 PATH 环境变量中查找
    import shutil
    for exe in exe_names:
        found = shutil.which(exe)
        if found:
            norm = os.path.normpath(found)
            if norm not in seen:
                seen.add(norm)
                results.append({"path": norm, "source": "PATH"})

    # 3. 查注册表（仅 Windows，非 Windows 静默跳过）
    try:
        import winreg
        if not hasattr(winreg, 'OpenKey'):
            raise ImportError("winreg not available")
        reg_paths = [
            r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths",
            r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths",
        ]
        for reg_path in reg_paths:
            for exe in exe_names:
                try:
                    key_path = os.path.join(reg_path, exe)
                    with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, key_path) as k:
                        val, _ = winreg.QueryValueEx(k, "")
                        norm = os.path.normpath(val.strip('"'))
                        if norm not in seen and os.path.isfile(norm):
                            seen.add(norm)
                            results.append({"path": norm, "source": "registry"})
                except (FileNotFoundError, OSError):
                    continue
    except ImportError:
        pass  # 非 Windows 环境

    if not results:
        return {"ok": False, "error": f"Program not found: {name}",
                "searched_names": exe_names}

    return {"ok": True, "name": name, "count": len(results), "results": results,
            "best": results[0]["path"]}


def _op_launch(params: dict) -> dict:
    """
    启动程序或打开文件。
    params:
      path  — 可执行文件完整路径（优先）
      name  — 程序名（自动查找后启动，需先有 find_program 结果或直接 name）
      args  — 启动参数列表（可选）
      wait  — 是否等待程序退出（默认 False）
    """
    import subprocess as _sp

    path = params.get("path", "").strip()
    name = params.get("name", "").strip()
    args = params.get("args", [])
    wait = params.get("wait", False)

    # 没有 path 时自动 find_program
    if not path and name:
        found = _op_find_program({"name": name})
        if not found.get("ok"):
            return {"ok": False, "error": f"Program not found '{name}': {found.get('error')}"}
        path = found["best"]

    if not path:
        return {"ok": False, "error": "path or name required"}

    if not os.path.isfile(path):
        # 尝试 shutil.which
        import shutil
        resolved = shutil.which(path)
        if resolved:
            path = resolved
        else:
            return {"ok": False, "error": f"Executable not found: {path}"}

    cmd = [path] + (args if isinstance(args, list) else [str(args)])

    try:
        if wait:
            proc = _sp.run(cmd, capture_output=True, text=True, timeout=60)
            return {"ok": proc.returncode == 0, "path": path,
                    "returncode": proc.returncode,
                    "stdout": proc.stdout.strip(), "stderr": proc.stderr.strip()}
        else:
            _sp.Popen(cmd, close_fds=True)
            return {"ok": True, "path": path, "status": "launched",
                    "message": f"Started: {os.path.basename(path)}"}
    except Exception as e:
        return {"ok": False, "error": f"Launch failed: {e}", "path": path}

_FILE_CHAT_URL = "http://127.0.0.1:4001/file-chat"
_FILE_CHAT_FALLBACK_URL = "http://127.0.0.1:3000/v1/file-chat"
_CHAT_URL = "http://127.0.0.1:3000/chat"

_BINARY_EXTS = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".ico", ".tiff",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv", ".pdf",
}


def _call_direct_chat(message: str, agent_id: str = "default") -> dict:
    """POST {message} to /chat — no file upload. Returns {ok, result, downloaded_b64, downloaded_ext}."""
    payload = json.dumps({
        "message": message,
        "agentId": agent_id,
        "newChat": False,
    }).encode("utf-8")
    try:
        req = urllib.request.Request(
            _CHAT_URL,
            data=payload,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=660) as resp:
            return json.loads(resp.read().decode("utf-8"))
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _call_file_chat(file_path: str, message: str = "", agent_id: str = "default") -> dict:
    """POST {file_path, message, agentId} to executor_server /file-chat → api-server → ChatGPT.
    Falls back directly to api-server :3000/v1/file-chat if executor_server is not running.
    """
    payload = json.dumps({
        "file_path": file_path,
        "message": message,
        "agentId": agent_id,
    }).encode("utf-8")

    for url in (_FILE_CHAT_URL, _FILE_CHAT_FALLBACK_URL):
        try:
            req = urllib.request.Request(
                url,
                data=payload,
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=660) as resp:
                return json.loads(resp.read().decode("utf-8"))
        except urllib.error.URLError:
            if url == _FILE_CHAT_FALLBACK_URL:
                return {"ok": False, "error": "Both executor_server (:4001) and api-server (:3000) are unreachable"}
            # executor_server not available — try api-server directly
            continue
        except Exception as e:
            return {"ok": False, "error": str(e)}


def _op_read_web(params: dict) -> dict:
    """
    将本地文件上传给 ChatGPT Web，返回 ChatGPT 的分析回复。
    支持任意文件类型（二进制、非明文等）。

    params:
      path      — 本地文件路径（必填）
      message   — 附带的提示语（可选，默认 "Please analyse this file."）
      agent_id  — ChatGPT 代理槽 ID（可选，默认 "default"）
    """
    path = _expand(params.get("path", ""))
    message = params.get("message", "Please analyse this file.")
    agent_id = params.get("agent_id", "default")

    if not path:
        return {"ok": False, "error": "path cannot be empty"}
    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}

    result = _call_file_chat(path, message, agent_id)
    if not result.get("ok"):
        return result
    return {
        "ok": True,
        "path": path,
        "reply": result.get("text", ""),
    }


def _op_write_web(params: dict) -> dict:
    """
    把本地文件上传给 ChatGPT，要求 ChatGPT 修改，再把结果写回本地（自动备份）。

    结果提取优先级：
      1. downloaded_b64   — ChatGPT 触发了真实文件下载（二进制直写，文本解码后写）
      2. 回复文本中含下载 URL — 用 urllib 抓取 URL 内容
      3. 回复文本中的代码围栏 — 剥除 ``` 取代码块
      4. 原始回复文本兜底（仅文本文件）

    params:
      path        — 本地源文件路径（必填）
      dst         — 写入目标路径（可选；不填则覆写 path）
      message     — 修改指令（必填）
      agent_id    — ChatGPT 代理槽 ID（默认 "default"）
      backup      — 是否备份原文件（默认 True）
      encoding    — 写回编码（默认 utf-8，仅文本）
      strip_fence — 是否剥除代码围栏（默认 True）
    """
    path = _expand(params.get("path", ""))
    dst  = _expand(params.get("dst", "")) or path   # write destination; default = overwrite source
    message = params.get("message", "")
    agent_id = params.get("agent_id", "default")
    do_backup = params.get("backup", True)
    encoding = params.get("encoding", "utf-8")
    strip_fence = params.get("strip_fence", True)

    if not path:
        return {"ok": False, "error": "path cannot be empty"}

    # New binary file → trigger direct generation (no file to upload)
    ext = os.path.splitext(path)[1].lower()
    if not os.path.isfile(path) and ext in _BINARY_EXTS:
        print(f"   [write_web] New binary file {path} — triggering direct generation via chat")
        gen_prompt = (
            f"IMPORTANT: Do NOT output any JSON block or file_op command. "
            f"Do NOT describe the output in text. "
            f"Just generate and display the result directly in this chat. "
            f"Task: {message}"
        )
        result = _call_direct_chat(gen_prompt, agent_id)
        if not result.get("ok"):
            return {"ok": False, "error": result.get("error", "Direct chat failed")}
        b64   = result.get("downloaded_b64", "")
        d_ext = result.get("downloaded_ext", ext or ".bin")
        if b64:
            import base64 as _b64mod
            dst_path = _expand(params.get("dst", "")) or path
            try:
                os.makedirs(os.path.dirname(dst_path) or ".", exist_ok=True)
                with open(dst_path, "wb") as _f:
                    _f.write(_b64mod.b64decode(b64))
                print(f"   [write_web] Saved generated file: {dst_path} ({os.path.getsize(dst_path)} bytes)")
                return {"ok": True, "bytes": os.path.getsize(dst_path), "path": dst_path,
                        "note": "Generated by ChatGPT via direct chat"}
            except Exception as e:
                return {"ok": False, "error": f"Failed to save generated file: {e}"}
        # No file captured — return the text reply as context
        return {"ok": False, "error": f"ChatGPT did not generate a downloadable file. Reply: {result.get('result','')[:200]}"}

    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}
    if not message:
        return {"ok": False, "error": "message (modification instruction) cannot be empty"}

    api_result = _call_file_chat(path, message, agent_id)
    if not api_result.get("ok"):
        return api_result

    raw_reply = api_result.get("text", "")
    downloaded_b64 = api_result.get("downloaded_b64", "")
    still_generating = api_result.get("generating", False)
    terminal_text_only = api_result.get("terminal_text_only", False)

    # If ChatGPT is still generating and there's no download content, report failure
    if still_generating and not downloaded_b64:
        return {
            "ok": False,
            "error": (
                "ChatGPT is still generating (timeout). "
                "The task may take longer than expected. "
                f"Partial reply: {raw_reply[:200]}"
            ),
        }
    if is_bin and terminal_text_only and not downloaded_b64:
        return {
            "ok": False,
            "error": (
                "ChatGPT accepted the upload but returned text only; "
                "no modified binary/image file was captured. "
                f"Reply: {raw_reply[:200]}"
            ),
        }

    raw_bytes = None   # binary result
    content = None     # text result
    source = ""

    # ── Priority 1: real file download captured by CDP ──────────────────────
    if downloaded_b64:
        try:
            raw_bytes = base64.b64decode(downloaded_b64)
            source = "cdp_download"
        except Exception:
            pass

    # ── Priority 2: download URL in reply text ───────────────────────────────
    if raw_bytes is None and content is None and raw_reply:
        url_match = re.search(
            r'https?://[^\s\)\]"\'<>]+\.(?:shader|glsl|hlsl|py|js|ts|cs|cpp|c|h|'
            r'json|txt|xml|yaml|yml|cfg|ini|zip|tar|gz|'
            r'png|jpg|jpeg|gif|bmp|webp)[^\s\)\]"\'<>]*',
            raw_reply, re.IGNORECASE
        )
        if url_match:
            dl_url = url_match.group(0)
            try:
                req = urllib.request.Request(dl_url, headers={"User-Agent": "Mozilla/5.0"})
                with urllib.request.urlopen(req, timeout=30) as resp:
                    raw_bytes = resp.read()
                source = f"url_download:{dl_url[:60]}"
            except Exception:
                pass

    # ── Priority 3: code fence in reply text (text files only) ──────────────
    if raw_bytes is None and content is None and raw_reply and strip_fence and not is_bin:
        fence_match = re.search(r"```[^\n]*\n([\s\S]*?)```", raw_reply, re.MULTILINE)
        if fence_match:
            content = fence_match.group(1)
            source = "code_fence"

    # ── Priority 4: raw reply as-is (text files only) ────────────────────────
    if raw_bytes is None and content is None and not is_bin:
        content = raw_reply
        source = "raw_reply"

    # ── Write back ───────────────────────────────────────────────────────────
    # Ensure destination directory exists
    Path(dst).parent.mkdir(parents=True, exist_ok=True)
    bak = _make_backup(path) if do_backup and os.path.isfile(path) else None

    if raw_bytes is not None:
        if not raw_bytes:
            return {"ok": False, "error": "ChatGPT returned empty file content", "path": dst,
                    "raw_reply": raw_reply[:200]}
        try:
            with open(dst, "wb") as f:
                f.write(raw_bytes)
            result = {"ok": True, "path": dst, "source": source, "bytes": len(raw_bytes),
                      "raw_reply_length": len(raw_reply)}
            if bak:
                result["backup"] = bak
            return result
        except Exception as e:
            return {"ok": False, "error": f"Binary write failed: {e}", "path": dst}
    else:
        if not content or not content.strip():
            return {"ok": False, "error": "ChatGPT returned empty content", "path": dst,
                    "raw_reply": raw_reply[:200]}
        try:
            _write_text(dst, content, encoding)
        except Exception as e:
            return {"ok": False, "error": f"Write failed: {e}", "path": dst}
        result = {"ok": True, "path": dst, "source": source,
                  "bytes": len(content.encode(encoding, errors="replace")),
                  "raw_reply_length": len(raw_reply)}
        if bak:
            result["backup"] = bak
        return result


def _op_backup(params: dict) -> dict:
    """手动备份文件（不修改原文件）。"""
    path = _expand(params.get("path", ""))
    if not os.path.isfile(path):
        return {"ok": False, "error": f"File not found: {path}"}
    bak = _make_backup(path)
    return {"ok": True, "path": path, "backup": bak}


def _op_history(params: dict) -> dict:
    """
    列出文件的历史备份。
    params:
      path — 目标文件路径
    """
    path = _expand(params.get("path", ""))
    if not path:
        return {"ok": False, "error": "path cannot be empty"}
    baks = _list_backups(path)
    return {"ok": True, "path": path, "count": len(baks), "backups": baks}


def _op_restore(params: dict) -> dict:
    """
    还原文件到某个历史版本。
    params:
      path    — 目标文件路径
      backup  — 备份文件完整路径（来自 history 的 backup 字段）；
                省略时还原到最近一次备份。
      backup  — 也可传 "latest"（默认）或 "oldest"
    还原前会先把当前版本再备份一次（保留还原路径）。
    """
    path = _expand(params.get("path", ""))
    bak_ref = params.get("backup", "latest")

    if not path:
        return {"ok": False, "error": "path cannot be empty"}

    baks = _list_backups(path)
    if not baks:
        return {"ok": False, "error": f"No backups found for: {path}"}

    if bak_ref in ("latest", ""):
        chosen = baks[0]          # list is newest-first
    elif bak_ref == "oldest":
        chosen = baks[-1]
    else:
        # explicit path
        expanded_ref = _expand(bak_ref)
        matched = [b for b in baks if b["backup"] == expanded_ref or b["name"] == bak_ref]
        if not matched:
            return {"ok": False, "error": f"Backup not found: {bak_ref}",
                    "available": [b["name"] for b in baks]}
        chosen = matched[0]

    bak_path = chosen["backup"]
    if not os.path.isfile(bak_path):
        return {"ok": False, "error": f"Backup file missing on disk: {bak_path}"}

    # Back up current state before restoring
    pre_restore_bak = _make_backup(path) if os.path.isfile(path) else None

    shutil.copy2(bak_path, path)
    result = {
        "ok": True, "path": path,
        "restored_from": bak_path,
        "restored_timestamp": chosen["modified"],
    }
    if pre_restore_bak:
        result["pre_restore_backup"] = pre_restore_bak
    return result


# ── 调度表 ────────────────────────────────────────────────

_ACTIONS = {
    "list":         _op_list,
    "find":         _op_find,
    "read":         _op_read,
    "write":        _op_write,
    "write_chunk":  _op_write_chunk,
    "append":       _op_append,
    "patch":        _op_patch,
    "insert":       _op_insert,
    "delete_lines": _op_delete_lines,
    "mkdir":        _op_mkdir,
    "copy":         _op_copy,
    "move":         _op_move,
    "delete":       _op_delete,
    "exists":       _op_exists,
    "stat":         _op_stat,
    "tree":         _op_tree,
    "find_program": _op_find_program,
    "launch":       _op_launch,
    "backup":       _op_backup,
    "restore":      _op_restore,
    "history":      _op_history,
    "read_web":     _op_read_web,
    "write_web":    _op_write_web,
}


# ── 公开入口 ──────────────────────────────────────────────

def run(action: str, params: dict) -> dict:
    """Execute one file operation; return dict with ok: bool."""
    fn = _ACTIONS.get(action)
    if fn is None:
        return {"ok": False,
                "error": f"Unknown action: {action}; supported: {list(_ACTIONS)}"}
    try:
        return fn(params)
    except Exception as e:
        return {"ok": False, "error": f"Error executing {action}: {e}"}


def schema_hint() -> str:
    """Return file_op format hint for AI. Load from prompts/file_ops_hint.txt if present."""
    _hint_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "prompts", "file_ops_hint.txt")
    if os.path.isfile(_hint_path):
        try:
            with open(_hint_path, "r", encoding="utf-8") as f:
                return f.read().strip()
        except Exception:
            pass
    return """
## file_op file operations

Use command: "file_op" for local files:

```json
{"command":"file_op","action":"<action>","path":"<path>","<key>":"<value>"}
```

| action       | Required             | Optional                          |
| list         | path                 |                                   |
| find         | path                 | pattern regex max_depth           |
| read         | path                 | line_start line_end               |
| write        | path content         | encoding backup                   |
| write_chunk  | path content         | line_start line_end encoding backup |
| append       | path content         | newline backup                    |
| patch        | path replacements    | use_regex backup                  |
| insert       | path line content    | after backup                      |
| delete_lines | path line_start      | line_end backup                   |
| mkdir        | path                 |                                   |
| copy         | src dst              |                                   |
| move         | src dst              |                                   |
| delete       | path                 | force                             |
| exists       | path                 |                                   |
| tree         | path                 | max_depth                         |
| find_program | name                 |                                   |
| launch       | path or name         | args wait                         |
| backup       | path                 |                                   |
| history      | path                 |                                   |
| restore      | path                 | backup (path/"latest"/"oldest")   |
| read_web     | path                 | message agent_id                  |
| write_web    | path message         | agent_id backup encoding          |

write_chunk replaces only lines line_start..line_end; other lines are preserved.
Use write_chunk for large files to avoid truncation issues.
backup/history/restore let you view and revert to previous versions.
read_web  — upload file to ChatGPT Web and return its analysis (supports any file type).
write_web — upload file to ChatGPT Web with a modification instruction, write the returned
            content back to dst (defaults to overwriting path). Use for binary/image files.
            Example (modify image and save to desktop):
            {"command":"file_op","action":"write_web","path":"C:\\source.png","dst":"C:\\Users\\admin\\Desktop\\output.png","message":"Refine the details"}
Paths support %USERPROFILE% %DESKTOP% etc.
"""
